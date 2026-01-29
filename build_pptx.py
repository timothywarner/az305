"""
AZ-305 Comprehensive Training Presentation Builder
Generates a 97-slide professional Azure-themed PPTX deck.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK_BLUE = RGBColor(0x00, 0x20, 0x50)
GREEN = RGBColor(0x00, 0xA3, 0x6C)
PINK = RGBColor(0xE8, 0x3E, 0x8C)
PURPLE = RGBColor(0x77, 0x19, 0xAA)
GOLD = RGBColor(0xFF, 0xB9, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
BLACK = RGBColor(0x00, 0x00, 0x00)
MID_GRAY = RGBColor(0x44, 0x44, 0x44)
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xFD)
LIGHT_GOLD = RGBColor(0xFF, 0xF8, 0xE1)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Segment labels for footer
SEGMENTS = {
    0: "Opening",
    1: "Segment 1: Identity, Governance & Monitoring",
    2: "Segment 2: Data Storage Solutions",
    3: "Segment 3: Business Continuity & HA",
    4: "Segment 4: Compute & Application Architecture",
    5: "Segment 5: Networking & Migrations",
    6: "Closing",
}

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Use blank layout
BLANK_LAYOUT = prs.slide_layouts[6]

slide_counter = 0


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------
def _set_fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_paragraph(text_frame, text, font_size=18, bold=False, color=BLACK,
                   alignment=PP_ALIGN.LEFT, space_before=Pt(4),
                   space_after=Pt(4), font_name="Segoe UI", level=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    return p


def add_footer(slide, segment_id, slide_num):
    """Add segment label bottom-left and slide number bottom-right."""
    seg_label = SEGMENTS.get(segment_id, "")
    _add_textbox(slide, Inches(0.5), Inches(7.0), Inches(6), Inches(0.4),
                 seg_label, font_size=10, color=MID_GRAY)
    _add_textbox(slide, Inches(12.0), Inches(7.0), Inches(1), Inches(0.4),
                 str(slide_num), font_size=10, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)


def add_header_bar(slide, title_text, subtitle_text=None):
    """Azure blue header bar across the top with title."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(0),
                                 SLIDE_WIDTH, Inches(1.2))
    _set_fill_solid(bar, AZURE_BLUE)
    bar.line.fill.background()

    _add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                 title_text, font_size=28, bold=True, color=WHITE,
                 font_name="Segoe UI Semibold")
    if subtitle_text:
        _add_textbox(slide, Inches(0.5), Inches(0.72), Inches(12), Inches(0.4),
                     subtitle_text, font_size=16, color=WHITE)


def new_content_slide(title, segment_id, subtitle=None):
    """Create a standard content slide with header bar and footer."""
    global slide_counter
    slide_counter += 1
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_header_bar(slide, title, subtitle)
    add_footer(slide, segment_id, slide_counter)
    return slide


def new_divider_slide(title, badge_text, segment_id):
    """Full dark-blue background divider slide."""
    global slide_counter
    slide_counter += 1
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0), Inches(0),
                                SLIDE_WIDTH, SLIDE_HEIGHT)
    _set_fill_solid(bg, DARK_BLUE)
    bg.line.fill.background()

    _add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(2),
                 title, font_size=44, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold")

    # Badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(4.5), Inches(4.5),
                                   Inches(4.3), Inches(0.8))
    _set_fill_solid(badge, GOLD)
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = badge_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = "Segoe UI Semibold"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(8)

    add_footer(slide, segment_id, slide_counter)
    return slide


def new_exam_tip_slide(title, segment_id):
    """Slide with light gold background for exam tips."""
    global slide_counter
    slide_counter += 1
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0), Inches(0),
                                SLIDE_WIDTH, SLIDE_HEIGHT)
    _set_fill_solid(bg, LIGHT_GOLD)
    bg.line.fill.background()

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(0),
                                 SLIDE_WIDTH, Inches(1.2))
    _set_fill_solid(bar, PINK)
    bar.line.fill.background()

    _add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                 title, font_size=28, bold=True, color=WHITE,
                 font_name="Segoe UI Semibold")
    _add_textbox(slide, Inches(0.5), Inches(0.72), Inches(12), Inches(0.4),
                 "EXAM TIPS", font_size=16, color=WHITE)

    add_footer(slide, segment_id, slide_counter)
    return slide


def add_bullet_list(slide, items, left=Inches(0.8), top=Inches(1.5),
                    width=Inches(11.5), height=Inches(5.2),
                    font_size=18, color=BLACK, bold_prefix=True):
    """Add a bulleted list to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        if bold_prefix and ": " in item:
            parts = item.split(": ", 1)
            run1 = p.add_run()
            run1.text = "\u2022 " + parts[0] + ": "
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = color
            run1.font.name = "Segoe UI"
            run2 = p.add_run()
            run2.text = parts[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = "Segoe UI"
        else:
            p.text = "\u2022 " + item
    return txBox


def add_table(slide, rows_data, col_widths, left=Inches(0.5),
              top=Inches(1.5), row_height=Inches(0.45),
              font_size=14, header_color=AZURE_BLUE):
    """Add a formatted table to a slide.
    rows_data: list of lists. First row is header."""
    num_rows = len(rows_data)
    num_cols = len(rows_data[0])
    total_width = sum(col_widths)

    table_shape = slide.shapes.add_table(num_rows, num_cols,
                                         left, top,
                                         Emu(int(total_width)),
                                         Emu(int(row_height * num_rows)))
    table = table_shape.table

    for ci, w in enumerate(col_widths):
        table.columns[ci].width = Emu(int(w))

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Segoe UI"
                if ri == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = BLACK
            cell.text_frame.paragraphs[0].space_before = Pt(2)
            cell.text_frame.paragraphs[0].space_after = Pt(2)

            if ri == 0:
                _set_cell_fill(cell, header_color)
            elif ri % 2 == 0:
                _set_cell_fill(cell, LIGHT_GRAY)
            else:
                _set_cell_fill(cell, WHITE)

    return table_shape


def _set_cell_fill(cell, color):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def add_review_questions(slide, questions):
    """Add numbered scenario questions."""
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5),
                                     Inches(11.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (q, a) in enumerate(questions):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(14)
        p.text = f"Q{i+1}: {q}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.font.name = "Segoe UI"

        pa = tf.add_paragraph()
        pa.text = f"   Answer: {a}"
        pa.font.size = Pt(15)
        pa.font.color.rgb = GREEN
        pa.font.name = "Segoe UI"
        pa.space_before = Pt(4)
        pa.space_after = Pt(4)


def add_demo_slide(slide, demos):
    """Format demo callout items."""
    # Icon-like shape
    icon = slide.shapes.add_shape(MSO_SHAPE.LIGHTNING_BOLT,
                                  Inches(0.8), Inches(1.5),
                                  Inches(0.6), Inches(0.6))
    _set_fill_solid(icon, GOLD)
    icon.line.fill.background()

    _add_textbox(slide, Inches(1.6), Inches(1.5), Inches(6), Inches(0.5),
                 "Live Demo Exercises", font_size=22, bold=True,
                 color=DARK_BLUE)

    add_bullet_list(slide, demos, left=Inches(1.0), top=Inches(2.3),
                    width=Inches(11), height=Inches(4.5), font_size=18)


# ============================================================================
# BUILD ALL 97 SLIDES
# ============================================================================

# ---------- SECTION 0: OPENING (Slides 1-8) ----------

# Slide 1: Title
slide_counter += 1
s = prs.slides.add_slide(BLANK_LAYOUT)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                        SLIDE_WIDTH, SLIDE_HEIGHT)
_set_fill_solid(bg, DARK_BLUE)
bg.line.fill.background()

# Accent bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.8),
                         SLIDE_WIDTH, Inches(0.06))
_set_fill_solid(bar, GOLD)
bar.line.fill.background()

_add_textbox(s, Inches(1), Inches(1.0), Inches(11), Inches(1.5),
             "AZ-305: Designing Microsoft Azure\nInfrastructure Solutions",
             font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER,
             font_name="Segoe UI Semibold")

_add_textbox(s, Inches(1), Inches(3.2), Inches(11), Inches(1),
             "Tim Warner  |  O'Reilly Live Learning  |  January 2026",
             font_size=22, color=GOLD, alignment=PP_ALIGN.CENTER)

_add_textbox(s, Inches(1), Inches(4.5), Inches(11), Inches(1),
             "Microsoft MVP  |  MCT  |  Azure Solutions Architect Expert",
             font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)

add_footer(s, 0, slide_counter)


# Slide 2: About the Instructor
s = new_content_slide("About the Instructor", 0)
add_bullet_list(s, [
    "Tim Warner: Microsoft MVP, Microsoft Certified Trainer (MCT)",
    "Certification: Azure Solutions Architect Expert (AZ-305)",
    "Platform: O'Reilly Live Learning instructor -- courses on Azure, DevOps, and cloud architecture",
    "Experience: 20+ years in IT training and consulting",
    "Contact: timothywarner316@gmail.com",
    "Social: @TechTrainerTim",
    "GitHub: github.com/timothywarner",
], font_size=20)


# Slide 3: Course Objectives
s = new_content_slide("Course Objectives", 0)
add_bullet_list(s, [
    "Segment 1: Design identity, governance, and monitoring solutions (25-30%)",
    "Segment 2: Design data storage solutions for relational, non-relational, and unstructured data (20-25%)",
    "Segment 3: Design business continuity, high availability, and disaster recovery solutions (15-20%)",
    "Segment 4: Design compute and application architecture solutions (~17%)",
    "Segment 5: Design networking solutions and plan migrations (~18%)",
], font_size=20)


# Slide 4: Exam Overview
s = new_content_slide("Exam Overview", 0, "AZ-305 Key Facts")
add_table(s, [
    ["Attribute", "Detail"],
    ["Exam Code", "AZ-305"],
    ["Title", "Designing Microsoft Azure Infrastructure Solutions"],
    ["Pass Score", "700 / 1000"],
    ["Duration", "120 minutes"],
    ["Questions", "~40-60 (multiple choice, case studies, drag-and-drop)"],
    ["Last Updated", "October 18, 2024"],
    ["Prerequisite", "AZ-104 recommended (not required)"],
    ["Cost", "$165 USD"],
], col_widths=[Inches(3), Inches(8)])


# Slide 5: Agenda & Schedule
s = new_content_slide("Agenda & Schedule", 0)
add_table(s, [
    ["Time", "Segment", "Topic", "Weight"],
    ["9:00 - 9:55", "Segment 1", "Identity, Governance & Monitoring", "25-30%"],
    ["10:05 - 10:55", "Segment 2", "Data Storage Solutions", "20-25%"],
    ["11:05 - 11:55", "Segment 3", "Business Continuity & HA", "15-20%"],
    ["12:05 - 12:55", "Segment 4", "Compute & Application Architecture", "~17%"],
    ["1:05 - 1:50", "Segment 5", "Networking & Migrations", "~18%"],
], col_widths=[Inches(2.2), Inches(2), Inches(5.5), Inches(1.8)],
    top=Inches(1.5))

_add_textbox(s, Inches(0.8), Inches(4.8), Inches(11), Inches(0.5),
             "10-minute breaks between each segment. All times Eastern.",
             font_size=16, color=MID_GRAY)


# Slide 6: Cross-Cutting Themes
s = new_content_slide("Cross-Cutting Themes", 0,
                      "Three pillars that appear in EVERY segment")
items = [
    "Well-Architected Framework: The lens through which every architecture decision is evaluated -- reliability, security, cost, operations, performance.",
    "Zero Trust: \"Never trust, always verify.\" Assume breach, verify explicitly, use least-privilege access. This model drives networking, identity, and data design.",
    "Managed Identity: Eliminate secrets from code. System-assigned and user-assigned identities replace connection strings, API keys, and passwords.",
]
add_bullet_list(s, items, font_size=18, top=Inches(1.6))


# Slide 7: Well-Architected Framework
s = new_content_slide("Azure Well-Architected Framework", 0, "5 Pillars")
add_table(s, [
    ["Pillar", "Focus", "Key Question"],
    ["Reliability", "Resiliency & recovery", "Can the system recover from failures?"],
    ["Security", "Threat protection", "How do we protect data and systems?"],
    ["Cost Optimization", "Manage costs", "Are we spending only what we need?"],
    ["Operational Excellence", "Operations processes", "Can we monitor and improve?"],
    ["Performance Efficiency", "Scaling & performance", "Does it meet demand efficiently?"],
], col_widths=[Inches(2.5), Inches(3.5), Inches(5.5)],
    top=Inches(1.5))


# Slide 8: How to Use This Session
s = new_content_slide("How to Use This Session", 0)
add_bullet_list(s, [
    "Decision Matrices: Map directly to exam question logic -- learn the selection criteria, not just the services.",
    "Architecture Diagrams: Visual patterns that appear in case studies -- understand component relationships.",
    "Exam Tips (pink slides): Specific test-day tactics -- common traps, default answers, and elimination strategies.",
    "Demo Callouts: Hands-on exercises you can replicate in your own Azure subscription or sandbox.",
    "Review Questions: Scenario-based practice at the end of each segment -- mirrors the actual exam format.",
    "Ask Questions: Use the chat throughout. Real-time Q&A helps everyone learn.",
], font_size=18)


# ---------- SECTION 1: IDENTITY, GOVERNANCE & MONITORING (Slides 9-27) ----------

# Slide 9: Divider
new_divider_slide("Identity, Governance\n& Monitoring", "25-30% of Exam", 1)

# Slide 10: Learning Objectives
s = new_content_slide("Segment 1 Learning Objectives", 1)
add_bullet_list(s, [
    "Design authentication and authorization solutions using Microsoft Entra ID",
    "Design governance solutions with management groups, policies, and RBAC",
    "Design monitoring and logging solutions with Azure Monitor, Sentinel, and Application Insights",
    "Design secure access to secrets, keys, and certificates using Key Vault with managed identities",
], font_size=20)

# Slide 11: Entra ID Core Concepts
s = new_content_slide("Microsoft Entra ID Core Concepts", 1)
add_bullet_list(s, [
    "Tenant: A dedicated instance of Entra ID representing an organization. Each tenant has a unique directory.",
    "Directory: The identity store within a tenant -- users, groups, applications, service principals.",
    "Subscription: A billing boundary linked to one tenant. Multiple subscriptions per tenant, one tenant per subscription.",
    "CRITICAL RENAME: 'Azure AD' is now 'Microsoft Entra ID'. The exam uses the NEW naming exclusively.",
    "Entra External ID: Replaces Azure AD B2C for consumer-facing identity scenarios.",
], font_size=18)

# Slide 12: Authentication Decision Matrix
s = new_content_slide("Authentication Decision Matrix", 1)
add_table(s, [
    ["Scenario", "Solution", "Key Detail"],
    ["Employees only", "Microsoft Entra ID", "Cloud-native, MFA, Conditional Access"],
    ["Partner organizations", "Entra External ID (B2B)", "Invite external users, cross-tenant access"],
    ["Consumer applications", "Entra External ID (B2C)", "Custom sign-up/sign-in, social logins"],
    ["Hybrid identity", "Entra Connect / Cloud Sync", "Sync on-prem AD to cloud"],
    ["Federated / complex", "AD FS", "On-prem federation, claims-based auth"],
    ["Multi-tenant SaaS", "Multi-tenant app registration", "Consent framework, service principals"],
], col_widths=[Inches(2.8), Inches(3.8), Inches(5)])

# Slide 13: Hybrid Identity Options
s = new_content_slide("Hybrid Identity Options", 1)
add_table(s, [
    ["Method", "How It Works", "Passwords in Cloud?", "Complexity", "Recommendation"],
    ["Password Hash Sync (PHS)", "Hash of hash synced to Entra ID", "Yes (hashed)", "Low", "DEFAULT -- start here"],
    ["Pass-Through Auth (PTA)", "Auth request forwarded to on-prem", "No", "Medium", "Regulatory requirement"],
    ["Federation (AD FS)", "Redirect to on-prem AD FS farm", "No", "High", "Only if required"],
], col_widths=[Inches(2.2), Inches(3.3), Inches(2.2), Inches(1.5), Inches(2.5)],
    font_size=13)

_add_textbox(s, Inches(0.8), Inches(4.5), Inches(11), Inches(1),
             "Exam Default: PHS is the recommended starting point. Cloud Sync is the modern lightweight alternative to Entra Connect.",
             font_size=16, bold=True, color=AZURE_BLUE)

# Slide 14: Conditional Access Architecture
s = new_content_slide("Conditional Access Architecture", 1,
                      "Policy = Assignments + Conditions + Controls")
add_bullet_list(s, [
    "Signals (Inputs): User/group, Cloud app, Location (named/IP), Device platform & state, Sign-in risk level (Identity Protection)",
    "Decisions (Grant Controls): Allow, Block, Require MFA, Require compliant device, Require Hybrid Entra joined, Require app protection policy",
    "Session Controls: App-enforced restrictions, Conditional Access App Control (MCAS), Sign-in frequency, Persistent browser session",
    "Key Principle: Policies are additive. If ANY policy blocks, access is denied. Most restrictive policy wins.",
    "Report-Only Mode: Test policies without enforcement. Always use this before enabling enforcement.",
    "Named Locations: Define trusted IPs/countries. Combine with policies to skip MFA on corporate networks.",
], font_size=17)

# Slide 15: Managed Identity Deep Dive
s = new_content_slide("Managed Identity Deep Dive", 1,
                      "Eliminate secrets from your code")
add_table(s, [
    ["Attribute", "System-Assigned", "User-Assigned"],
    ["Lifecycle", "Tied to resource (deleted together)", "Independent (manage separately)"],
    ["Sharing", "1:1 with resource", "1:Many across resources"],
    ["Use Case", "Single resource needs access", "Multiple resources share same identity"],
    ["Creation", "Enable on the resource", "Create as standalone resource, then assign"],
    ["Example", "VM accessing Key Vault", "Multiple VMs accessing same storage account"],
    ["Exam Default", "Start with system-assigned", "Use when sharing is needed"],
], col_widths=[Inches(2), Inches(4.8), Inches(4.8)],
    font_size=13)

# Slide 16: RBAC & Authorization
s = new_content_slide("RBAC & Authorization", 1)
add_bullet_list(s, [
    "Scope Hierarchy: Management Group > Subscription > Resource Group > Resource. Permissions inherit downward.",
    "Built-in Roles: Owner (full + assign), Contributor (full - assign), Reader (view only), User Access Admin (assign only).",
    "Custom Roles: JSON definition with Actions, NotActions, DataActions, NotDataActions. Assignable at MG, Sub, or RG scope.",
    "Deny Assignments: Override RBAC allows. Created by Blueprints and managed apps. Cannot create directly.",
    "Best Practice: Assign roles to GROUPS, not individual users. Use PIM for just-in-time elevation.",
    "Scope Tip: Assign at the NARROWEST scope possible. If access is needed at RG level, do not assign at subscription level.",
], font_size=17)

# Slide 17: PIM & Identity Governance
s = new_content_slide("PIM & Identity Governance", 1)
add_bullet_list(s, [
    "Privileged Identity Management (PIM): Just-in-time (JIT) role activation. Eligible vs Active assignments.",
    "PIM Workflow: User requests activation -> Approval (optional) -> Time-limited access (e.g., 8 hours) -> Auto-expires",
    "Access Reviews: Periodic review of who has access to what. Self-review, manager review, or group owner review.",
    "Entitlement Management: Access packages bundle resources (groups, apps, SharePoint). Catalog + Policies + Requests.",
    "Identity Governance: Lifecycle workflows for joiner/mover/leaver. Automate onboarding and offboarding.",
    "Exam Key: PIM requires Entra ID P2 license. Access Reviews also require P2.",
], font_size=17)

# Slide 18: Management Group Hierarchy
s = new_content_slide("Management Group Hierarchy", 1,
                      "Cloud Adoption Framework Landing Zone")
add_bullet_list(s, [
    "Root MG: Tenant Root Group (auto-created, do NOT apply restrictive policies here)",
    "  Platform MG: Shared services -- contains Identity, Management, and Connectivity subscriptions",
    "    Identity MG: Domain controllers, Entra Connect servers",
    "    Management MG: Log Analytics, Automation, Monitor resources",
    "    Connectivity MG: Hub VNet, VPN/ExpressRoute gateways, DNS, Firewall",
    "  Landing Zones MG: Workload subscriptions",
    "    Corp MG: Internal apps connected to corporate network",
    "    Online MG: Internet-facing apps",
    "  Sandbox MG: Experimentation subscriptions (no connectivity to corp)",
    "  Decommissioned MG: Subscriptions being retired",
    "Depth Limit: 6 levels deep (excluding Root and subscription level)",
], font_size=16)

# Slide 19: Azure Policy Deep Dive
s = new_content_slide("Azure Policy Deep Dive", 1)
add_table(s, [
    ["Effect", "What It Does", "When to Use"],
    ["Deny", "Blocks non-compliant resource creation/update", "Enforce hard requirements"],
    ["Audit", "Logs non-compliance, allows creation", "Monitor before enforcing"],
    ["Append", "Adds fields to resource during creation", "Add tags, IP rules"],
    ["DeployIfNotExists", "Deploys related resource if missing", "Auto-configure diagnostics"],
    ["Modify", "Changes properties on existing resources", "Add/update tags, settings"],
    ["AuditIfNotExists", "Audits if related resource is missing", "Check for diagnostics settings"],
    ["Disabled", "Policy exists but not enforced", "Testing, temporary disable"],
], col_widths=[Inches(2.5), Inches(4.5), Inches(4.5)],
    font_size=13)

_add_textbox(s, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
             "Initiative = Collection of policies assigned together. Use initiatives for compliance standards (e.g., CIS, NIST).",
             font_size=15, bold=True, color=AZURE_BLUE)

# Slide 20: Governance Decision Matrix
s = new_content_slide("Governance Decision Matrix", 1)
add_table(s, [
    ["Requirement", "Tool", "Details"],
    ["Enforce naming conventions", "Azure Policy (Deny)", "Regex patterns on resource names"],
    ["Control costs", "Budgets + Cost Management + Policy", "Alerts at thresholds, deny expensive SKUs"],
    ["Regulatory compliance", "Policy Initiatives", "Built-in: CIS, NIST, ISO, PCI DSS"],
    ["Resource organization", "Tags + Management Groups", "Cost center, environment, owner tags"],
    ["Prevent accidental deletion", "Resource Locks", "CanNotDelete or ReadOnly locks"],
    ["Standardize deployments", "Blueprints / Template Specs", "Versioned, governed ARM/Bicep templates"],
    ["Track changes", "Activity Log + Change Analysis", "Who changed what and when"],
], col_widths=[Inches(2.8), Inches(3.5), Inches(5.2)],
    font_size=13)

# Slide 21: Azure Monitor Architecture
s = new_content_slide("Azure Monitor Architecture", 1)
add_bullet_list(s, [
    "Data Sources: Applications (App Insights), OS (agents), Azure resources (diagnostics), Subscriptions (activity log), Tenant (Entra logs)",
    "Collection: Data Collection Rules (DCR) and Diagnostic Settings route data to destinations",
    "Destinations: Log Analytics workspace (KQL queries), Metrics (near real-time), Storage (archive), Event Hub (streaming)",
    "Analysis: KQL queries, Workbooks (visual dashboards), Alerts (action groups: email, SMS, webhook, Logic App, Function)",
    "Azure Monitor Agent (AMA): Replaces legacy agents (MMA, Telegraf). Single agent for Windows and Linux. Uses DCRs.",
    "Log Analytics Retention: 30 days free interactive, up to 730 days, archive up to 12 years.",
], font_size=17)

# Slide 22: Application Insights & Observability
s = new_content_slide("Application Insights & Observability", 1)
add_bullet_list(s, [
    "APM: Application Performance Management -- automatic instrumentation for .NET, Java, Node.js, Python",
    "Distributed Tracing: End-to-end transaction tracking across microservices. Correlates with Operation ID.",
    "Live Metrics: Real-time stream of requests, failures, dependencies. Zero-cost when viewing.",
    "Availability Tests: URL ping (global), Standard test (status + SSL + content), Custom TrackAvailability.",
    "Smart Detection: AI-driven anomaly detection for failures, performance degradation, memory leaks.",
    "Application Map: Visual topology of dependencies. Shows call rates, latency, failure rates between components.",
    "Workspace-based: Always use workspace-based App Insights (not classic). Required for cross-resource queries.",
], font_size=17)

# Slide 23: Microsoft Sentinel
s = new_content_slide("Microsoft Sentinel", 1, "Cloud-Native SIEM + SOAR")
add_bullet_list(s, [
    "SIEM: Security Information and Event Management -- collect, detect, investigate, respond",
    "SOAR: Security Orchestration, Automation and Response -- playbooks (Logic Apps) for automated response",
    "Data Connectors: 100+ built-in. Microsoft 365, Entra ID, Defender, AWS, GCP, firewalls, custom (CEF/Syslog)",
    "Analytics Rules: Scheduled (KQL), Microsoft Security (from Defender), Fusion (ML correlation), NRT (near real-time)",
    "Workbooks: Interactive dashboards for SOC analysts. Built-in templates for common scenarios.",
    "When Sentinel vs Defender for Cloud: Sentinel = full SIEM for SOC teams. Defender = CSPM + workload protection. Use both together.",
], font_size=17)

# Slide 24: Key Vault Patterns
s = new_content_slide("Key Vault Patterns", 1)
add_bullet_list(s, [
    "Access Model: RBAC (recommended) vs Vault Access Policies (legacy). RBAC integrates with Entra ID roles.",
    "Soft Delete: Enabled by default, 7-90 day retention. CANNOT be disabled on new vaults.",
    "Purge Protection: When enabled, soft-deleted items cannot be force-purged during retention period. Enable for production.",
    "Managed Identity Access: Grant Key Vault Secrets User role to managed identity. Zero secrets in code.",
    "Secret Rotation: Event Grid notifies when secret nears expiry -> Function rotates -> Stores new secret in Key Vault.",
    "Key Types: RSA & EC keys, software or HSM-backed. Premium SKU for HSM. Managed HSM for FIPS 140-2 L3.",
    "Network Security: Private Endpoint for VNet access. Firewall rules for IP restrictions. Disable public access.",
], font_size=16)

# Slide 25: Demo Callout
s = new_content_slide("Segment 1 Demos", 1)
add_demo_slide(s, [
    "Conditional Access Policy: Create a policy requiring MFA for Azure portal access from non-corporate locations",
    "Azure Policy Assignment: Assign 'Require tag on resource group' policy with Deny effect to a test subscription",
    "Managed Identity + Key Vault: Create a VM with system-assigned identity, grant Key Vault access, retrieve secret from code",
    "Log Analytics KQL: Query sign-in logs to find failed authentications and risky sign-ins",
])

# Slide 26: Exam Tips
s = new_exam_tip_slide("Segment 1 Exam Tips", 1)
add_bullet_list(s, [
    "PHS is the DEFAULT hybrid identity recommendation. Choose PTA only if passwords must never leave on-prem.",
    "Managed Identity eliminates secrets. If a question mentions connection strings or API keys, Managed Identity is likely the answer.",
    "Policy > Locks for enforcement. Locks prevent deletion; Policies prevent non-compliant creation.",
    "PIM provides JIT access. If a question asks about 'least privilege for admins', PIM is the answer.",
    "Sentinel = SIEM. If the question mentions 'correlate security events across sources', think Sentinel.",
    "Log Analytics retention: 30 days free interactive. Archive tier for long-term at lower cost.",
], font_size=18, top=Inches(1.5), color=DARK_BLUE)

# Slide 27: Review Questions
s = new_content_slide("Segment 1 Review Questions", 1)
add_review_questions(s, [
    ("Your company requires that IT admins only have elevated privileges when performing admin tasks, with approval from a security manager. What should you implement?",
     "PIM (Privileged Identity Management) with eligible assignments and approval workflow. Admins request activation, security manager approves, access auto-expires."),
    ("A web application needs to read secrets from Key Vault without storing any credentials in code or configuration. The app runs on Azure App Service. What do you recommend?",
     "Enable system-assigned managed identity on the App Service. Grant it the 'Key Vault Secrets User' RBAC role on the Key Vault. Use DefaultAzureCredential in code."),
    ("You need to enforce that all resources in a subscription must have a 'CostCenter' tag. Resources without the tag should be blocked from creation. What do you use?",
     "Azure Policy with the 'Require a tag on resources' built-in definition, set with Deny effect. Assign at subscription scope."),
])


# ---------- SECTION 2: DATA STORAGE (Slides 28-45) ----------

# Slide 28: Divider
new_divider_slide("Data Storage Solutions", "20-25% of Exam", 2)

# Slide 29: Learning Objectives
s = new_content_slide("Segment 2 Learning Objectives", 2)
add_bullet_list(s, [
    "Design storage solutions for relational data (SQL Database, SQL MI, SQL on VM)",
    "Design for semi-structured and NoSQL data (Cosmos DB, Table Storage)",
    "Design for unstructured data (Blob Storage, ADLS Gen2, Azure Files)",
    "Design data integration and analytics architectures (ADF, Synapse, Event Hubs)",
], font_size=20)

# Slide 30: Storage Decision Tree
s = new_content_slide("Storage Decision Tree", 2)
add_table(s, [
    ["Data Type", "Service", "When to Use"],
    ["Structured / Relational", "Azure SQL Database", "Cloud-native, auto-scaling, managed"],
    ["Structured / High compat", "Azure SQL Managed Instance", "SQL Server compat, VNet native"],
    ["Structured / Full control", "SQL Server on Azure VM", "OS-level access, legacy features"],
    ["Semi-structured / Global", "Azure Cosmos DB", "Multi-model, global dist, low latency"],
    ["Semi-structured / Simple", "Table Storage", "Key-value, low cost, simple queries"],
    ["Unstructured / Objects", "Blob Storage", "Images, documents, backups, media"],
    ["Unstructured / Analytics", "ADLS Gen2", "Big data, hierarchical namespace, Spark"],
    ["File shares", "Azure Files", "SMB/NFS, lift-and-shift, AD DS integration"],
    ["High-perf file shares", "Azure NetApp Files", "SAP, HPC, Oracle, low-latency NFS"],
], col_widths=[Inches(2.5), Inches(3.5), Inches(5.5)],
    font_size=13, top=Inches(1.4))

# Slide 31: Storage Redundancy
s = new_content_slide("Storage Redundancy Matrix", 2)
add_table(s, [
    ["Option", "Copies", "Scope", "Use Case", "Read Secondary?"],
    ["LRS", "3", "Single datacenter", "Dev/test, non-critical", "No"],
    ["ZRS", "3", "3 availability zones", "Production, zone resilience", "No"],
    ["GRS", "6", "2 regions (primary+secondary)", "DR, cross-region", "No (failover only)"],
    ["GZRS", "6", "3 zones + secondary region", "Mission-critical", "No (failover only)"],
    ["RA-GRS", "6", "2 regions", "DR + read from secondary", "Yes (read-only)"],
    ["RA-GZRS", "6", "3 zones + secondary region", "Maximum resilience + read", "Yes (read-only)"],
], col_widths=[Inches(1.5), Inches(1), Inches(3), Inches(3.5), Inches(2.5)],
    font_size=13)

# Slide 32: Blob Tiers
s = new_content_slide("Blob Storage Tiers & Lifecycle", 2)
add_table(s, [
    ["Tier", "Access", "Min Duration", "Retrieval Cost", "Storage Cost"],
    ["Hot", "Frequent", "None", "Low", "Highest"],
    ["Cool", "Infrequent (30+ days)", "30 days", "Medium", "Lower"],
    ["Cold", "Rarely (90+ days)", "90 days", "Higher", "Lower still"],
    ["Archive", "Offline (180+ days)", "180 days", "Highest (rehydrate)", "Lowest"],
], col_widths=[Inches(1.5), Inches(2.8), Inches(2), Inches(2.5), Inches(2.5)],
    font_size=14)
_add_textbox(s, Inches(0.8), Inches(4.5), Inches(11), Inches(1.5),
             "Lifecycle Management: Automate tier transitions with rules.\n"
             "Example: Move to Cool after 30 days, Archive after 90, delete after 365.\n"
             "Rehydration: Archive requires rehydration (Standard: up to 15 hours, High priority: <1 hour).",
             font_size=16, color=DARK_BLUE)

# Slide 33: Azure Files & NetApp
s = new_content_slide("Azure Files & NetApp Files", 2)
add_bullet_list(s, [
    "Azure Files: Managed SMB (445) and NFS (2049) file shares in the cloud",
    "Tiers: Premium (SSD, low latency), Transaction Optimized (HDD), Hot, Cool",
    "AD DS Integration: Entra Domain Services or on-prem AD DS for identity-based access with NTFS permissions",
    "Azure File Sync: Sync on-prem file servers with Azure Files. Cloud tiering frees local disk space. Multi-site sync.",
    "Azure NetApp Files: Enterprise NAS -- sub-millisecond latency, Oracle, SAP HANA, HPC workloads",
    "NetApp tiers: Standard, Premium, Ultra. Capacity pools with volume allocation.",
    "Exam tip: Azure Files for general file shares; NetApp Files for high-performance or SAP/Oracle.",
], font_size=17)

# Slide 34: Relational DB Decision Matrix
s = new_content_slide("Relational Database Decision Matrix", 2)
add_table(s, [
    ["Feature", "SQL Database", "SQL Managed Instance", "SQL on VM"],
    ["Managed?", "Fully managed PaaS", "Fully managed PaaS", "IaaS (you manage OS)"],
    ["SQL Compat", "~95% (some gaps)", "~99% (near full)", "100% (full engine)"],
    ["VNet Native", "No (PE/service endpoint)", "Yes (in VNet)", "Yes (in VNet)"],
    ["SQL Agent", "Elastic Jobs", "Yes", "Yes"],
    ["Cross-DB Queries", "Elastic Query", "Yes", "Yes"],
    ["CLR / Linked Servers", "No", "Yes", "Yes"],
    ["Best For", "New cloud apps", "Migration / compat", "Legacy / full control"],
    ["Cost Model", "DTU or vCore", "vCore only", "VM + license"],
], col_widths=[Inches(2.3), Inches(3), Inches(3), Inches(3.2)],
    font_size=12, top=Inches(1.4))

# Slide 35: Azure SQL Architecture Patterns
s = new_content_slide("Azure SQL Architecture Patterns", 2)
add_bullet_list(s, [
    "Elastic Pools: Share resources across multiple databases. Ideal for multi-tenant SaaS with variable workloads.",
    "Hyperscale: Up to 100 TB, rapid scale-out reads (up to 4 read replicas), fast backups regardless of size.",
    "Serverless: Auto-pause after idle period, auto-scale compute. Pay only for compute used. Best for intermittent workloads.",
    "DTU vs vCore: DTU = bundled (CPU+IO+memory). vCore = choose independently. vCore recommended for new deployments.",
    "Azure Hybrid Benefit: Use existing SQL Server licenses for ~55% savings. Apply at server or pool level.",
    "Ledger tables: Tamper-evident tables with blockchain-like verification. Regulatory and audit scenarios.",
], font_size=17)

# Slide 36: Cosmos DB Design
s = new_content_slide("Cosmos DB Design", 2, "Global Distribution & Multi-Model")
add_bullet_list(s, [
    "APIs: NoSQL (native, recommended), MongoDB, Cassandra, Gremlin (graph), Table, PostgreSQL",
    "Global Distribution: Multi-region writes for active-active. Single-region writes + multi-region reads for active-passive.",
    "Partition Key Rules: (1) High cardinality -- many distinct values. (2) Even distribution -- no hot partitions. (3) Included in queries -- avoid cross-partition queries.",
    "RU/s: Request Units per second. Provisioned (predictable) or Autoscale (variable, min 10% of max). Serverless for dev/test.",
    "Item Size: Max 2 MB per item. Design for smaller documents with proper denormalization.",
    "Change Feed: Stream of changes for event-driven patterns. Trigger Functions, update caches, sync data.",
], font_size=17)

# Slide 37: Cosmos DB Consistency
s = new_content_slide("Cosmos DB Consistency Spectrum", 2)
add_table(s, [
    ["Level", "Guarantee", "Latency", "Throughput", "Use Case"],
    ["Strong", "Linearizable reads", "Highest", "Lowest", "Financial transactions"],
    ["Bounded Staleness", "Reads lag by k versions or t time", "High", "Lower", "Leaderboards, counters"],
    ["Session (DEFAULT)", "Read-your-own-writes in session", "Medium", "Medium", "Most applications"],
    ["Consistent Prefix", "Reads never see out-of-order writes", "Lower", "Higher", "Social updates"],
    ["Eventual", "No ordering guarantee", "Lowest", "Highest", "Counters, likes, non-critical"],
], col_widths=[Inches(2), Inches(3), Inches(1.5), Inches(1.5), Inches(3.5)],
    font_size=13)
_add_textbox(s, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
             "Exam Default: Session consistency is the default and suits most applications.",
             font_size=16, bold=True, color=AZURE_BLUE)

# Slide 38: Data Platform Architecture
s = new_content_slide("Data Platform Architecture", 2, "Medallion Pattern")
add_bullet_list(s, [
    "Bronze Layer (Raw): Ingest raw data as-is. ADLS Gen2 with hierarchical namespace. Parquet or Delta format.",
    "Silver Layer (Cleansed): Validated, deduplicated, conformed. Schema enforcement. Business logic applied.",
    "Gold Layer (Curated): Aggregated, business-ready. Star schema for BI. Served via Synapse SQL pools or Power BI.",
    "ADLS Gen2: Storage foundation. Hierarchical namespace + Blob APIs. ACLs for fine-grained access.",
    "Synapse Analytics: Unified analytics -- serverless SQL, dedicated SQL pools, Spark, Pipelines, Data Explorer.",
    "Microsoft Fabric: Next-gen unified analytics platform. Integrates Power BI, Data Factory, Synapse, Real-Time Analytics.",
], font_size=17)

# Slide 39: Data Integration
s = new_content_slide("Data Integration", 2)
add_table(s, [
    ["Service", "Pattern", "Use Case"],
    ["Azure Data Factory", "ETL / ELT orchestration", "Batch data movement, 90+ connectors, mapping data flows"],
    ["Synapse Pipelines", "Same ADF engine", "Data integration within Synapse workspace"],
    ["Event Hubs", "Streaming ingestion", "Millions of events/sec, Kafka compatible, capture to storage"],
    ["Stream Analytics", "Real-time processing", "SQL-like queries on streams, windowing functions"],
    ["Azure Databricks", "Spark-based analytics", "ML, data engineering, collaborative notebooks"],
    ["Logic Apps", "Integration workflows", "B2B, SaaS connectors, low-code orchestration"],
], col_widths=[Inches(2.5), Inches(3), Inches(6)])

# Slide 40: Data Protection & Encryption
s = new_content_slide("Data Protection & Encryption", 2)
add_bullet_list(s, [
    "At Rest: Storage Service Encryption (SSE) for blobs/files. TDE for SQL databases. Always on by default.",
    "In Transit: TLS 1.2+ enforced. Minimum TLS version configurable on storage accounts.",
    "In Use: Always Encrypted for SQL (client-side encryption). Confidential computing for VMs.",
    "Key Management: Microsoft-managed keys (MMK, default) or Customer-managed keys (CMK, Key Vault).",
    "CMK Benefits: You control rotation, you can revoke access. Required by some compliance standards.",
    "Double Encryption: Infrastructure encryption adds second layer. Enable at storage account creation.",
], font_size=17)

# Slide 41: Private Endpoints for Data
s = new_content_slide("Private Endpoints for Data", 2,
                      "Zero Trust Networking for Data Services")
add_bullet_list(s, [
    "Architecture: VNet -> Private Endpoint (NIC with private IP) -> Private DNS Zone -> Data Service",
    "DNS Resolution: privatelink.blob.core.windows.net resolves to PE private IP, not public IP",
    "Disable Public Access: After PE is configured, disable public network access on the service",
    "On-prem Access: Conditional DNS forwarder forwards privatelink.* queries to Azure DNS (168.63.129.16)",
    "Supported Services: Storage, SQL, Cosmos DB, Key Vault, App Config, Event Hubs, Service Bus, and 60+ more",
    "Cost: ~$7.30/month per PE + data processing charges. Small cost for significant security improvement.",
], font_size=17)

# Slide 42: Data Residency & Purview
s = new_content_slide("Data Residency & Microsoft Purview", 2)
add_bullet_list(s, [
    "Sovereign Regions: Azure Government (US), Azure China (21Vianet). Data stays within sovereign boundary.",
    "Data Residency: Choose regions that meet regulatory requirements (GDPR, LGPD, etc.). Paired regions for DR.",
    "Microsoft Purview: Unified data governance -- data catalog, lineage tracking, classifications, access policies.",
    "Data Catalog: Scan and classify data across Azure, on-prem, and multi-cloud. Automatic sensitivity labels.",
    "Data Lineage: Visual tracking of data flow from source to consumption. Supports ADF, Synapse, Power BI.",
    "Exam Context: Know when Purview is the answer -- 'discover, classify, govern data estate' keywords.",
], font_size=17)

# Slide 43: Demo
s = new_content_slide("Segment 2 Demos", 2)
add_demo_slide(s, [
    "Storage Lifecycle: Configure blob lifecycle management to transition from Hot to Cool to Archive",
    "Cosmos DB Configuration: Create a Cosmos DB account with NoSQL API, select partition key, set consistency",
    "Private Endpoint for SQL: Create a PE for Azure SQL Database with private DNS zone integration",
    "Data Factory Pipeline: Build a simple copy pipeline from Blob Storage to SQL Database",
])

# Slide 44: Exam Tips
s = new_exam_tip_slide("Segment 2 Exam Tips", 2)
add_bullet_list(s, [
    "ZRS for zone resilience in production. LRS only for dev/test or non-critical data.",
    "SQL Managed Instance is the migration sweet spot -- near 100% SQL Server compatibility with full PaaS management.",
    "Session consistency is the Cosmos DB default. Strong only when absolutely required (higher latency, lower throughput).",
    "Private Endpoints = Zero Trust for data. If the question mentions 'secure access' to a data service, think PE.",
    "Archive tier = offline. You CANNOT read archived blobs without rehydrating first. Plan for up to 15 hours.",
    "Partition key rules: High cardinality, even distribution, included in WHERE clause of queries.",
], font_size=18, top=Inches(1.5), color=DARK_BLUE)

# Slide 45: Review
s = new_content_slide("Segment 2 Review Questions", 2)
add_review_questions(s, [
    ("A global e-commerce app needs sub-10ms read latency worldwide. Data is semi-structured product catalog. Users mostly read their own recent changes. Which database and consistency?",
     "Cosmos DB with NoSQL API, multi-region reads, Session consistency (default). Partition key = productCategory or productId depending on query patterns."),
    ("You are migrating a SQL Server 2019 database that uses SQL Agent jobs, cross-database queries, and CLR assemblies. The team wants full PaaS management. What do you recommend?",
     "Azure SQL Managed Instance. It supports SQL Agent, cross-DB queries, and CLR -- features not available in SQL Database. Fully managed PaaS."),
    ("A healthcare company must store patient images for 7 years at minimum cost, with rare access. Images must never leave the US East region. What storage design?",
     "Blob Storage in US East, LRS or ZRS. Lifecycle policy: move to Archive tier after 30 days. RA-GRS NOT needed (single region requirement). Use immutability policies for retention."),
])


# ---------- SECTION 3: BCDR (Slides 46-60) ----------

new_divider_slide("Business Continuity\n& High Availability", "15-20% of Exam", 3)

# Slide 47
s = new_content_slide("Segment 3 Learning Objectives", 3)
add_bullet_list(s, [
    "Design for high availability using availability zones, sets, and load balancing",
    "Design backup and recovery solutions with Azure Backup and site recovery",
    "Design disaster recovery strategies for multi-region architectures",
    "Calculate composite SLAs and map RTO/RPO to Azure services",
], font_size=20)

# Slide 48
s = new_content_slide("RTO vs RPO Fundamentals", 3)
add_bullet_list(s, [
    "RPO (Recovery Point Objective): Maximum acceptable DATA LOSS measured in time. How much data can you afford to lose?",
    "RTO (Recovery Time Objective): Maximum acceptable DOWNTIME measured in time. How quickly must you recover?",
    "Example: RPO = 1 hour means you can lose up to 1 hour of data. RTO = 4 hours means you must recover within 4 hours.",
    "Cost Relationship: Lower RPO/RTO = higher cost. Near-zero RPO requires synchronous replication. Near-zero RTO requires hot standby.",
    "Business drives the numbers: Finance and stakeholders define acceptable RPO/RTO, architects design to meet them.",
    "Exam Pattern: Scenario gives RPO/RTO requirements -> you select the appropriate Azure service/configuration.",
], font_size=17)

# Slide 49
s = new_content_slide("RTO / RPO Decision Matrix", 3)
add_table(s, [
    ["RPO Requirement", "Solution", "RTO Requirement", "Solution"],
    ["RPO = 0 (zero loss)", "Sync replication / AZ", "RTO < 1 min", "Auto-failover (SQL FG, Cosmos)"],
    ["RPO < 15 min", "Continuous replication / ASR", "RTO < 1 hour", "Hot standby / ASR failover"],
    ["RPO < 1 hour", "Frequent backups / log shipping", "RTO < 4 hours", "Warm standby / ASR"],
    ["RPO < 24 hours", "Daily backups", "RTO < 24 hours", "Cold / backup restore"],
    ["RPO = days", "Weekly backups / geo-restore", "RTO = days", "Rebuild from backup"],
], col_widths=[Inches(2.5), Inches(3.5), Inches(2.5), Inches(3)])

# Slide 50
s = new_content_slide("Availability Zones vs Availability Sets", 3)
add_table(s, [
    ["Attribute", "Availability Zones", "Availability Sets"],
    ["Protection Level", "Datacenter failure", "Rack/hardware failure"],
    ["SLA", "99.99%", "99.95%"],
    ["Spread", "3 physically separate datacenters", "Up to 3 fault + 20 update domains"],
    ["Network Latency", "< 2 ms between zones", "N/A (same datacenter)"],
    ["Cost", "Cross-zone data transfer charged", "No extra cost"],
    ["Example", "VMs in Zone 1, 2, 3 + zone LB", "VMs in FD 0, 1, 2 behind LB"],
    ["Recommendation", "DEFAULT for production", "Legacy or regions without zones"],
], col_widths=[Inches(2.5), Inches(4.5), Inches(4.5)],
    font_size=13)

# Slide 51
s = new_content_slide("SLA Calculation", 3)
add_bullet_list(s, [
    "Single VM SLAs: Premium SSD = 99.9%, Availability Set = 99.95%, Availability Zones = 99.99%",
    "Composite SLA (serial): Multiply individual SLAs. Web (99.95%) x DB (99.99%) = 99.94%",
    "Composite SLA (parallel/redundant): 1 - (1-SLA_A) x (1-SLA_B). Two 99.9% instances = 99.9999%",
    "Example: Web tier (2 VMs in zones, 99.99%) x App tier (99.95%) x SQL (99.99%) = 99.93%",
    "Impact: 99.9% = ~8.76 hrs/year downtime. 99.99% = ~52 min/year. 99.95% = ~4.38 hrs/year.",
    "Strategy: Add redundancy to the WEAKEST component. That is where you get the most SLA improvement.",
], font_size=17)

# Slide 52
s = new_content_slide("Azure Backup Architecture", 3)
add_bullet_list(s, [
    "Recovery Services Vault: Central management for backup and ASR. Supports VMs, SQL, Files, SAP, on-prem.",
    "Backup Vault: Newer vault type for Azure Disks, Blobs, PostgreSQL, AKS. Simpler management.",
    "Agents: VM Extension (Azure VMs), MARS (files/folders to cloud), MABS/DPM (on-prem workloads).",
    "Policies: Daily/weekly/monthly/yearly retention. Customize per workload. Instant restore from snapshots.",
    "Storage Redundancy: LRS (default), ZRS (zone protection), GRS (cross-region). Choose at vault creation.",
    "Soft Delete: 14 additional days to recover deleted backup data. Enhanced soft delete for ransomware protection.",
    "Cross-Region Restore: Enable on GRS vaults to restore in paired region even when primary is healthy.",
], font_size=16)

# Slide 53
s = new_content_slide("Azure Site Recovery (ASR)", 3)
add_bullet_list(s, [
    "Azure-to-Azure: Replicate VMs between Azure regions. RPO ~30 seconds for VMs.",
    "On-prem-to-Azure: VMware VMs, Hyper-V VMs, physical servers. Process Server handles replication.",
    "Test Failover: Failover to isolated VNet. NO impact to production. Validate DR plan safely.",
    "Planned Failover: Zero data loss. Shut down source, replicate final changes, bring up target.",
    "Unplanned Failover: Source is down. Use latest recovery point. Some data loss possible (within RPO).",
    "Recovery Plans: Group VMs, define startup order, add scripts/manual actions. Automate entire DR sequence.",
    "Key Distinction: ASR = DR (disaster recovery), NOT backup. Use Azure Backup for backup.",
], font_size=17)

# Slide 54
s = new_content_slide("SQL BCDR Options", 3)
add_table(s, [
    ["Option", "RPO", "Scope", "Auto Failover?", "Best For"],
    ["Zone Redundant", "0", "Same region, 3 zones", "Automatic", "HA within region"],
    ["Active Geo-Replication", "~5 seconds", "Up to 4 secondaries, any region", "Manual", "Read scale, custom DR"],
    ["Auto-Failover Groups", "~5 seconds", "1 secondary region", "Automatic (DNS)", "Production DR"],
    ["Geo-Restore", "~1 hour", "Paired region", "Manual (restore)", "Budget DR"],
    ["PITR", "5 min (up to 35 days)", "Same region", "Manual (restore)", "Accidental changes"],
], col_widths=[Inches(2.2), Inches(1.3), Inches(3.5), Inches(2), Inches(2.5)],
    font_size=13)

# Slide 55
s = new_content_slide("Multi-Region Patterns", 3)
add_table(s, [
    ["Pattern", "RTO", "RPO", "Cost", "Complexity", "Use Case"],
    ["Active-Active", "~0", "~0", "Highest", "Highest", "Global apps, zero downtime"],
    ["Active-Passive (Hot)", "Minutes", "Near-zero", "High", "Medium", "Critical apps"],
    ["Active-Passive (Warm)", "Minutes-hours", "Minutes", "Medium", "Medium", "Important apps"],
    ["Active-Passive (Cold)", "Hours", "Hours", "Low", "Low", "Non-critical"],
], col_widths=[Inches(2.2), Inches(1.3), Inches(1.5), Inches(1.5), Inches(1.5), Inches(3.5)],
    font_size=14)
_add_textbox(s, Inches(0.8), Inches(4.5), Inches(11), Inches(1),
             "Traffic Routing: Azure Front Door (global HTTP) or Traffic Manager (DNS-based) for multi-region failover.",
             font_size=16, bold=True, color=AZURE_BLUE)

# Slide 56
s = new_content_slide("HA Decision Matrix by Service", 3)
add_table(s, [
    ["Service", "HA Mechanism", "SLA", "Key Configuration"],
    ["VMs", "Availability Zones + Load Balancer", "99.99%", "Zone-redundant deployment"],
    ["Azure SQL", "Auto-Failover Groups", "99.99%", "Automatic DNS failover"],
    ["Cosmos DB", "Multi-region writes", "99.999%", "5-nines with multi-write"],
    ["App Service", "Zone-redundant plan", "99.99%", "Premium v3 + zone enabled"],
    ["AKS", "Zone-spanning node pools", "99.99%", "3 zones, system + user pools"],
    ["Storage", "ZRS / GZRS", "99.99%+", "Zone or geo-zone redundant"],
    ["Functions", "Zone-redundant Premium", "99.99%", "Premium plan required"],
], col_widths=[Inches(1.8), Inches(3.2), Inches(1.3), Inches(5.2)],
    font_size=13)

# Slide 57
s = new_content_slide("DR Testing", 3)
add_bullet_list(s, [
    "Test Failover (ASR): Failover VMs to isolated VNet. No impact to production. Validate apps work.",
    "Test Failover (SQL): Failover groups support test failover with alternate DNS suffix.",
    "Recovery Plans: Define multi-VM failover with sequencing. Add pre/post scripts for configuration.",
    "Documented Runbooks: Step-by-step DR procedures. Who does what, in what order, with what approvals.",
    "Regular DR Drills: Quarterly or semi-annual. Test RTO/RPO claims. Document lessons learned.",
    "Chaos Engineering: Azure Chaos Studio for fault injection. Test resilience under controlled failure.",
], font_size=17)

# Slide 58: Demo
s = new_content_slide("Segment 3 Demos", 3)
add_demo_slide(s, [
    "Azure Backup Policy: Create a Recovery Services Vault with daily backup policy for Azure VMs",
    "ASR Replication: Enable Azure-to-Azure replication for a VM, view replication health",
    "SQL Failover Group: Create an auto-failover group between two Azure SQL servers in different regions",
    "Zone Deployment: Deploy a VM to a specific availability zone and verify zone placement",
])

# Slide 59: Exam Tips
s = new_exam_tip_slide("Segment 3 Exam Tips", 3)
add_bullet_list(s, [
    "Auto-failover groups > active geo-replication for SQL DR. Groups provide automatic DNS failover.",
    "ASR is for DR, NOT backup. If the question asks about 'recovering deleted files', use Azure Backup.",
    "Availability Zones = 99.99% SLA. This is the default HA recommendation for production workloads.",
    "Front Door for global HTTP HA. Traffic Manager for DNS-based global routing (non-HTTP or simple).",
    "Test failover = no production impact. Always recommend test failover for DR validation.",
    "Composite SLA math: Multiply serial components, use parallel formula for redundant components.",
], font_size=18, top=Inches(1.5), color=DARK_BLUE)

# Slide 60: Review
s = new_content_slide("Segment 3 Review Questions", 3)
add_review_questions(s, [
    ("Your SQL database must survive a complete Azure region outage with automatic failover and near-zero RPO. Users connect via a single DNS name. What do you implement?",
     "Auto-failover group across two regions. Provides automatic DNS failover (~5 second RPO), single read-write endpoint, and optional read-only endpoint."),
    ("A web application runs on 3 VMs behind a load balancer. The SLA for each VM is 99.9%. The LB SLA is 99.99%. What is the composite SLA for the web tier?",
     "Parallel VMs: 1 - (1-0.999)^3 = 99.9999%. Serial with LB: 99.9999% x 99.99% = 99.9899%. The LB becomes the limiting factor."),
    ("Your company needs to protect Azure VMs with RPO < 30 minutes and RTO < 1 hour for DR to a secondary region. Which service?",
     "Azure Site Recovery (ASR). Provides ~30 second RPO for Azure VMs with continuous replication, and minutes RTO with automated failover via recovery plans."),
])


# ---------- SECTION 4: COMPUTE & APP ARCH (Slides 61-76) ----------

new_divider_slide("Compute & Application\nArchitecture", "~17% of Exam", 4)

# Slide 62
s = new_content_slide("Segment 4 Learning Objectives", 4)
add_bullet_list(s, [
    "Select appropriate compute solutions (VMs, containers, serverless, PaaS)",
    "Design container strategies with AKS, Container Apps, and ACI",
    "Design serverless solutions with Azure Functions and Logic Apps",
    "Design application architecture patterns including messaging, caching, and API management",
], font_size=20)

# Slide 63
s = new_content_slide("Compute Decision Tree", 4)
add_table(s, [
    ["Requirement", "Service", "Key Differentiator"],
    ["Full OS control", "VMs / VMSS", "Custom software, legacy apps, Windows/Linux"],
    ["Container orchestration", "AKS", "Full Kubernetes, complex microservices"],
    ["Serverless containers", "Container Apps", "RECOMMENDED DEFAULT. Dapr, KEDA, simple scaling"],
    ["Event-driven code", "Azure Functions", "Per-execution billing, triggers + bindings"],
    ["Web applications", "App Service", "Managed PaaS, deployment slots, custom domains"],
    ["Simple container run", "Azure Container Instances", "Quick burst, sidecar, no orchestration"],
    ["Batch / HPC", "Azure Batch", "Large-scale parallel, job scheduling"],
    ["Desktop apps", "Azure Virtual Desktop", "VDI, multi-session Windows, RemoteApp"],
], col_widths=[Inches(2.5), Inches(3), Inches(6)])

# Slide 64
s = new_content_slide("VM Design Patterns", 4)
add_bullet_list(s, [
    "Series Selection: B (burstable, dev/test), D (general), E (memory, SAP/DB), N (GPU, AI/ML), L (storage, big data), F (compute-optimized)",
    "VMSS: Virtual Machine Scale Sets. Autoscale 0-1000 instances. Uniform (identical VMs) or Flexible (mixed) mode.",
    "Proximity Placement Groups: Co-locate VMs in same datacenter for lowest latency. Used with E/N series for SAP/HPC.",
    "Dedicated Hosts: Physical server dedicated to your org. Compliance, licensing (BYOL), isolated hardware.",
    "Spot VMs: Up to 90% discount, can be evicted. Batch jobs, dev/test, stateless workloads. NOT for production SLAs.",
    "Ephemeral OS Disks: Use local VM storage for OS disk. Faster reimaging, lower latency, no storage cost. Stateless VMs only.",
], font_size=17)

# Slide 65
s = new_content_slide("App Service Architecture", 4)
add_bullet_list(s, [
    "Plans: Free/Shared (dev) -> Basic (dedicated, no scaling) -> Standard (autoscale, slots) -> Premium (more perf, VNet integration) -> Isolated (ASE, full VNet)",
    "Deployment Slots: Swap staging to production with zero downtime. Slot-specific settings (connection strings, app settings).",
    "VNet Integration: Regional VNet integration (outbound from app to VNet). Private Endpoint (inbound to app from VNet).",
    "Managed Certificates: Free TLS certificates for custom domains in Standard+ plans. Auto-renewal.",
    "Health Checks: Built-in health check path. Removes unhealthy instances from LB rotation. Minimum 2 instances.",
    "App Service Environment (ASE) v3: Single-tenant, fully in your VNet. Internal LB for private-only apps. Zone redundant.",
], font_size=17)

# Slide 66
s = new_content_slide("Container Strategy", 4)
add_table(s, [
    ["Attribute", "AKS", "Container Apps", "ACI"],
    ["Orchestrator", "Full Kubernetes", "Managed (Kubernetes under hood)", "None"],
    ["Scaling", "Cluster + pod autoscaler", "KEDA (event-driven, 0 to N)", "Manual (1:1)"],
    ["Networking", "Full K8s networking, CNI", "Built-in ingress, Envoy", "VNet or public IP"],
    ["Service Mesh", "Istio, Linkerd, OSM", "Dapr (built-in)", "N/A"],
    ["Complexity", "High (K8s expertise needed)", "Low-Medium", "Lowest"],
    ["Cost", "Node VMs + management", "Per-request/vCPU+memory", "Per-second vCPU+memory"],
    ["Best For", "Complex microservices", "Most containerized apps", "Simple tasks, burst"],
    ["RECOMMENDATION", "When K8s is mandated", "DEFAULT for containers", "Sidecar / quick run"],
], col_widths=[Inches(2), Inches(3.5), Inches(3.5), Inches(2.5)],
    font_size=12, top=Inches(1.4))

# Slide 67
s = new_content_slide("AKS Deep Dive", 4)
add_bullet_list(s, [
    "Node Pools: System pool (CoreDNS, kube-proxy) + User pools (workloads). Different VM sizes per pool.",
    "Cluster Autoscaler: Scale nodes 0-100+ based on pending pods. Combine with Horizontal Pod Autoscaler.",
    "Workload Identity: Replace pod identity with federated OIDC tokens. Managed identity bound to K8s service account.",
    "Ingress: NGINX, Application Gateway Ingress Controller (AGIC), or managed NGINX. TLS termination at ingress.",
    "Service Mesh: Istio add-on (managed by Azure). mTLS, traffic management, observability.",
    "When AKS is Warranted: Team has K8s expertise, need for custom operators, complex networking, multi-cluster.",
], font_size=17)

# Slide 68
s = new_content_slide("Container Apps with Dapr", 4)
add_bullet_list(s, [
    "Dapr Service Invocation: Service-to-service calls with built-in retries, mTLS, and observability.",
    "Dapr Pub/Sub: Decouple services with topics. Supports Service Bus, Event Hubs, Redis as brokers.",
    "Dapr State Management: Pluggable state stores (Cosmos DB, Redis, SQL). Consistent API across stores.",
    "Dapr Bindings: Input/output bindings to external services (queues, databases, HTTP endpoints).",
    "KEDA Scaling: Scale to zero and back. Scale based on queue depth, HTTP requests, CPU, custom metrics.",
    "Revisions & Traffic Splitting: Blue-green and canary deployments. Split traffic between revisions by percentage.",
    "Jobs: Run-to-completion tasks, scheduled (cron) or event-triggered. Parallel execution.",
], font_size=17)

# Slide 69
s = new_content_slide("Serverless: Azure Functions", 4)
add_table(s, [
    ["Plan", "Scale", "Max Timeout", "VNet", "Min Instances", "Cost Model"],
    ["Consumption", "0 to 200 instances", "10 min (default)", "No", "0", "Per execution"],
    ["Premium", "10-100 instances", "Unlimited", "Yes", "1 (pre-warmed)", "vCPU/memory/sec"],
    ["Dedicated (App Svc)", "Manual/autoscale", "Unlimited", "Yes", "Plan-defined", "App Service plan"],
], col_widths=[Inches(2), Inches(2.2), Inches(2), Inches(1), Inches(1.8), Inches(2.5)],
    font_size=13, top=Inches(1.5))
add_bullet_list(s, [
    "Durable Functions: Stateful workflows. Patterns: Chaining, Fan-out/Fan-in, Async HTTP, Monitor, Human interaction.",
    "Triggers: HTTP, Timer, Queue, Blob, Event Grid, Event Hub, Cosmos DB change feed, Service Bus.",
    "Bindings: Input (read data) and Output (write data) without SDK code. Declarative in function.json.",
    "Cold Start: Consumption plan has cold start latency. Use Premium plan or Flex Consumption to avoid.",
], top=Inches(4.2), font_size=16)

# Slide 70
s = new_content_slide("Messaging Service Selection", 4)
add_table(s, [
    ["Pattern", "Service", "Key Feature", "Example"],
    ["Commands / Transactions", "Service Bus", "Queues + Topics, sessions, FIFO, dead-letter", "Order processing"],
    ["Events / Notifications", "Event Grid", "Push delivery, filtering, serverless events", "Blob created, resource changed"],
    ["Streaming / Telemetry", "Event Hubs", "Millions/sec, partitions, Kafka compatible", "IoT telemetry, clickstream"],
    ["Simple queue", "Storage Queue", "64KB messages, millions queued, low cost", "Background jobs"],
], col_widths=[Inches(2.3), Inches(2), Inches(4), Inches(3.2)])
_add_textbox(s, Inches(0.8), Inches(4.5), Inches(11), Inches(1.5),
             "Exam Rule of Thumb:\n"
             "  'Commands' or 'exactly once' = Service Bus\n"
             "  'Events' or 'react to changes' = Event Grid\n"
             "  'Stream' or 'millions of events' = Event Hubs",
             font_size=16, color=DARK_BLUE)

# Slide 71
s = new_content_slide("Application Architecture Patterns", 4)
add_bullet_list(s, [
    "Microservices: Independent services with own data stores. Deploy, scale, update independently. Use with Container Apps or AKS.",
    "Event-Driven: Services react to events asynchronously. Event Grid for discrete events, Event Hubs for streaming.",
    "CQRS: Command Query Responsibility Segregation. Separate read and write models. Cosmos DB + SQL read replicas.",
    "Saga Pattern: Distributed transactions across services. Choreography (events) or Orchestration (coordinator).",
    "Queue-Based Load Leveling: Buffer requests with Service Bus queue. Backend processes at its own pace. Prevents overload.",
    "Retry + Circuit Breaker: Transient fault handling. Exponential backoff. Circuit breaker prevents cascading failures.",
], font_size=17)

# Slide 72
s = new_content_slide("API Management (APIM)", 4)
add_table(s, [
    ["Tier", "Gateway", "Developer Portal", "VNet", "Multi-region", "Use Case"],
    ["Consumption", "Serverless", "No", "No", "No", "Low traffic APIs"],
    ["Developer", "Shared", "Yes", "No", "No", "Dev/test"],
    ["Basic", "Dedicated", "Yes", "No", "No", "Small production"],
    ["Standard", "Dedicated", "Yes", "External", "No", "Production"],
    ["Premium", "Dedicated", "Yes", "Internal + External", "Yes", "Enterprise"],
], col_widths=[Inches(1.5), Inches(1.5), Inches(2), Inches(2), Inches(1.5), Inches(3)],
    font_size=13, top=Inches(1.5))
add_bullet_list(s, [
    "Policies: Rate limiting (throttling), JWT validation, caching, request/response transformation, IP filtering.",
    "Internal Mode: APIM inside VNet with private IP only. Use App Gateway as frontend for public access.",
], top=Inches(4.8), font_size=16)

# Slide 73
s = new_content_slide("Caching Strategy", 4)
add_table(s, [
    ["Tier", "Memory", "SLA", "Features"],
    ["Basic", "250MB - 53GB", "No SLA", "Dev/test only"],
    ["Standard", "250MB - 53GB", "99.9%", "Replication, production"],
    ["Premium", "6GB - 120GB per shard", "99.9%", "Clustering, persistence, VNet, geo-replication"],
    ["Enterprise", "12GB - 2TB", "99.99%", "RediSearch, RedisBloom, RedisTimeSeries, active geo-replication"],
], col_widths=[Inches(1.5), Inches(2.5), Inches(1.3), Inches(6.2)],
    font_size=13, top=Inches(1.5))
add_bullet_list(s, [
    "Cache-Aside Pattern: App checks cache first. On miss, reads from DB, writes to cache. Most common pattern.",
    "Session Store: Store session state in Redis instead of in-memory. Enables stateless app tier with sticky sessions off.",
    "Output Caching: Cache rendered pages or API responses. Reduce backend load for read-heavy workloads.",
], top=Inches(4.3), font_size=16)

# Slide 74: Demo
s = new_content_slide("Segment 4 Demos", 4)
add_demo_slide(s, [
    "Container Apps + Dapr: Deploy a microservice app with Dapr pub/sub and KEDA scaling to zero",
    "Functions + Service Bus: Create an Azure Function triggered by a Service Bus queue message",
    "APIM Policies: Configure rate limiting and JWT validation policies on an API",
    "App Service Deployment Slots: Deploy to staging slot, test, then swap to production with zero downtime",
])

# Slide 75: Exam Tips
s = new_exam_tip_slide("Segment 4 Exam Tips", 4)
add_bullet_list(s, [
    "Container Apps is the DEFAULT container recommendation. Choose AKS only when full Kubernetes control is needed.",
    "AKS when: Custom operators, advanced networking (CNI), multi-cluster, team has K8s expertise.",
    "Service Bus for commands (ordered, transactional). Event Grid for events (reactive, push). Do not confuse them.",
    "Event Hubs for streaming (high volume, partitioned, Kafka). Not for request/response patterns.",
    "APIM for API governance: rate limiting, versioning, developer portal. Internal mode for private APIs.",
    "Durable Functions for stateful serverless: function chaining, fan-out/fan-in, human interaction patterns.",
], font_size=18, top=Inches(1.5), color=DARK_BLUE)

# Slide 76: Review
s = new_content_slide("Segment 4 Review Questions", 4)
add_review_questions(s, [
    ("Your team wants to deploy containerized microservices with event-driven autoscaling (including scale to zero) and built-in service-to-service communication. They do not have Kubernetes expertise. What do you recommend?",
     "Azure Container Apps with Dapr enabled and KEDA scaling rules. Provides event-driven scaling, service invocation, pub/sub, and state management without Kubernetes complexity."),
    ("An order processing system must guarantee FIFO (first-in-first-out) message delivery with exactly-once processing and dead-letter support. Which messaging service?",
     "Azure Service Bus with sessions enabled for FIFO ordering. Service Bus provides exactly-once delivery, dead-letter queues, and session-based ordering. Storage Queue does NOT support FIFO."),
    ("You need to expose internal APIs to external partners with rate limiting, API key authentication, and a developer portal. The APIs must only be accessible through the gateway. What do you design?",
     "Azure API Management (Standard or Premium tier). Configure rate-limit policies, subscription keys, and developer portal. Premium with internal VNet mode if APIs must be fully private."),
])


# ---------- SECTION 5: NETWORKING & MIGRATIONS (Slides 77-92) ----------

new_divider_slide("Networking\n& Migrations", "~18% of Exam", 5)

# Slide 78
s = new_content_slide("Segment 5 Learning Objectives", 5)
add_bullet_list(s, [
    "Design network solutions including VNets, hub-spoke, and Virtual WAN",
    "Design hybrid connectivity with VPN Gateway and ExpressRoute",
    "Design network security with Azure Firewall, NSGs, and DDoS Protection",
    "Plan and design migration strategies using Azure Migrate and CAF",
], font_size=20)

# Slide 79
s = new_content_slide("VNet Architecture", 5)
add_bullet_list(s, [
    "Address Space Planning: Use RFC 1918 ranges (10.0.0.0/8, 172.16.0.0/12, 192.168.0.0/16). Plan for growth.",
    "Non-Overlapping Addresses: Critical for peering and VPN. Hub and spokes must have unique CIDR ranges.",
    "Subnets: Segment by workload type. Some services require dedicated subnets (AKS, App Service VNet Integration, Bastion, Firewall).",
    "NSGs: Network Security Groups. Layer 4 (IP, port, protocol) rules. Attach to subnet or NIC. Processed by priority (100-4096).",
    "ASGs: Application Security Groups. Group VMs logically (web, app, db). Use in NSG rules instead of IP addresses.",
    "Service Endpoints vs Private Endpoints: SE = optimized route, still public IP. PE = private IP in your VNet. PE is preferred (Zero Trust).",
], font_size=17)

# Slide 80
s = new_content_slide("Hub-Spoke Topology", 5)
add_bullet_list(s, [
    "Hub VNet: Centralized shared services -- Azure Firewall, VPN Gateway, ExpressRoute Gateway, Bastion, DNS.",
    "Spoke VNets: Workload isolation. Each spoke is a separate VNet peered to the hub.",
    "Peering: VNet peering is non-transitive. Spoke A cannot reach Spoke B through the hub without UDR + Firewall/NVA.",
    "UDRs (User-Defined Routes): Force spoke traffic through the hub Firewall. 0.0.0.0/0 -> Firewall private IP.",
    "Gateway Transit: Hub's VPN/ER gateway shared with spokes. Enable 'Allow gateway transit' on hub, 'Use remote gateways' on spoke.",
    "DNS: Azure Private DNS Zones linked to hub. Spoke VNets use hub's DNS or Azure DNS Private Resolver.",
    "Scale: Suitable for up to ~10-15 spokes. Beyond that, consider Virtual WAN.",
], font_size=17)

# Slide 81
s = new_content_slide("Azure Virtual WAN", 5)
add_bullet_list(s, [
    "Managed Hub-Spoke: Microsoft manages the hub router, VPN/ER gateways, routing tables.",
    "When to Use: 10+ branch offices, global transit (any-to-any), SD-WAN integration, multiple ER circuits.",
    "Hub: Virtual WAN hub is a managed VNet. You cannot deploy custom VMs inside it (except NVAs from marketplace).",
    "Routing: Automatic any-to-any connectivity between spokes, branches, and ExpressRoute. Route tables for segmentation.",
    "Secured Virtual Hub: Integrated Azure Firewall Manager for centralized security policies.",
    "vs Manual Hub-Spoke: Virtual WAN = less control, more automation, better for scale. Manual = more control, better for customization.",
], font_size=17)

# Slide 82
s = new_content_slide("Hybrid Connectivity Decision", 5)
add_table(s, [
    ["Attribute", "VPN Gateway", "ExpressRoute"],
    ["Connection", "Over public internet (encrypted)", "Private connection (NOT encrypted by default)"],
    ["Bandwidth", "Up to 10 Gbps (VpnGw5)", "50 Mbps to 100 Gbps (Direct)"],
    ["Latency", "Variable (internet dependent)", "Predictable, low latency"],
    ["Redundancy", "Active-active with 2 tunnels", "Dual circuits to different peering locations"],
    ["Cost", "Lower (gateway + data egress)", "Higher (circuit + gateway + provider fees)"],
    ["Setup Time", "Minutes to hours", "Days to weeks (provider circuit)"],
    ["Use Case", "Dev/test, backup path, small offices", "Production, latency-sensitive, large data"],
    ["Max Resilience", "VPN + ExpressRoute together", "2 ER circuits + VPN as backup"],
], col_widths=[Inches(2), Inches(5), Inches(4.5)],
    font_size=13)

# Slide 83
s = new_content_slide("ExpressRoute Deep Dive", 5)
add_bullet_list(s, [
    "Peering Types: Azure Private Peering (VNets), Microsoft Peering (M365, Dynamics, Azure PaaS public IPs).",
    "ExpressRoute Direct: Dedicated 10/100 Gbps port pair at peering location. Massive data transfer, MACsec encryption.",
    "Global Reach: Connect on-prem sites through Microsoft backbone. Site A (ER) <-> Microsoft <-> Site B (ER). No VPN needed.",
    "FastPath: Bypass the ExpressRoute gateway for data-plane traffic. Ultra performance for latency-sensitive workloads.",
    "Encryption: ExpressRoute is NOT encrypted by default. Use MACsec (Direct) or IPsec VPN over ER for encryption.",
    "Billing: Unlimited data (flat rate) or Metered data (per GB egress). Premium add-on for global route access.",
], font_size=17)

# Slide 84
s = new_content_slide("Network Security Layers", 5, "Defense in Depth")
add_bullet_list(s, [
    "Layer 1 - DDoS Protection: Standard plan for VNet resources. Adaptive tuning, cost protection, rapid response.",
    "Layer 2 - Azure Front Door WAF: Global HTTP protection. OWASP rules, bot protection, geo-filtering, rate limiting.",
    "Layer 3 - Azure Firewall: Centralized egress and east-west filtering. FQDN rules, threat intelligence, TLS inspection.",
    "Layer 4 - NSGs: Distributed L4 filtering at subnet and NIC. Allow/deny by IP, port, protocol. Flow logs for visibility.",
    "Layer 5 - Application Identity: Managed identity, RBAC, service endpoints/private endpoints. No network path = no attack surface.",
    "Layer 6 - Encryption: TLS in transit, SSE/TDE at rest, Always Encrypted in use. End-to-end protection.",
], font_size=17)

# Slide 85
s = new_content_slide("Firewall vs NSG vs NVA", 5)
add_table(s, [
    ["Attribute", "NSG", "Azure Firewall", "NVA (3rd Party)"],
    ["OSI Layer", "L3-L4", "L3-L7", "L3-L7 (vendor dependent)"],
    ["Cost", "Free", "~$1.25/hr + data processing", "VM cost + license"],
    ["Management", "Distributed (per subnet/NIC)", "Centralized in hub", "Self-managed or vendor console"],
    ["FQDN Filtering", "No", "Yes (application rules)", "Yes (vendor dependent)"],
    ["TLS Inspection", "No", "Yes (Premium)", "Yes (most vendors)"],
    ["Threat Intelligence", "No", "Yes (alert/deny known bad IPs)", "Yes (vendor feeds)"],
    ["Best For", "Micro-segmentation", "Centralized egress control", "Specific vendor requirements"],
], col_widths=[Inches(2), Inches(3.2), Inches(3.5), Inches(2.8)],
    font_size=12, top=Inches(1.4))

# Slide 86
s = new_content_slide("Private Link Architecture", 5)
add_bullet_list(s, [
    "Private Endpoint: NIC with a private IP from your VNet, connected to an Azure service via Private Link.",
    "Private DNS Zone: Maps service FQDN (e.g., mydb.database.windows.net) to the PE private IP.",
    "DNS Resolution Flow: Client -> Azure DNS -> privatelink zone -> PE private IP -> Service (via Microsoft backbone).",
    "Disable Public Access: Once PE is configured, disable public network access on the service for full Zero Trust.",
    "On-Prem Access: DNS conditional forwarder for privatelink.* zones -> Azure DNS (168.63.129.16) via VPN/ER.",
    "Private Link Service: Expose YOUR services behind a Standard LB as Private Link endpoints for consumers in other VNets/tenants.",
], font_size=17)

# Slide 87
s = new_content_slide("Load Balancing Decision", 5)
add_table(s, [
    ["Scope", "HTTP(S)?", "Service", "Key Features"],
    ["Global", "Yes", "Azure Front Door", "CDN, WAF, SSL offload, path routing, caching"],
    ["Global", "No", "Traffic Manager", "DNS-based, health probes, priority/weighted/geographic"],
    ["Regional", "Yes", "Application Gateway", "L7 LB, WAF, SSL termination, URL routing, rewrite"],
    ["Regional", "No", "Azure Load Balancer", "L4 LB, HA Ports, cross-region (Global tier)"],
], col_widths=[Inches(1.3), Inches(1.3), Inches(3), Inches(5.9)],
    font_size=14)
_add_textbox(s, Inches(0.8), Inches(4.5), Inches(11), Inches(1.5),
             "Common Combo: Front Door (global) -> Application Gateway (regional WAF) -> Load Balancer (backend VMs).\n"
             "Exam Shortcut: 'Global + HTTP' = Front Door. 'Regional + HTTP' = App Gateway. 'Non-HTTP' = LB or Traffic Manager.",
             font_size=16, color=DARK_BLUE)

# Slide 88
s = new_content_slide("Migration Strategies (5 Rs)", 5)
add_table(s, [
    ["Strategy", "Definition", "Effort", "Azure Service"],
    ["Rehost (Lift & Shift)", "Move as-is to Azure VMs", "Low", "Azure Migrate, ASR"],
    ["Replatform", "Minor optimizations during migration", "Low-Medium", "SQL MI, App Service"],
    ["Refactor", "Re-architect for cloud-native", "High", "Container Apps, Functions"],
    ["Rebuild", "Rewrite from scratch", "Highest", "New cloud-native services"],
    ["Replace", "Switch to SaaS product", "Varies", "M365, Dynamics, SaaS apps"],
], col_widths=[Inches(2.5), Inches(3.5), Inches(1.5), Inches(4)])
_add_textbox(s, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
             "CAF Migration Phases: Discover -> Assess -> Plan -> Migrate -> Optimize -> Govern -> Manage",
             font_size=16, bold=True, color=AZURE_BLUE)

# Slide 89
s = new_content_slide("Migration Tools", 5)
add_table(s, [
    ["Tool", "Purpose", "Key Capability"],
    ["Azure Migrate", "Discovery + assessment + migration", "VM, DB, web app, VDI assessment. Agentless or agent-based."],
    ["Database Migration Service", "Database migration", "SQL Server, PostgreSQL, MySQL to Azure. Online (minimal downtime) + Offline."],
    ["Azure Data Box", "Offline large data transfer", "100TB (Data Box) or 1PB (Data Box Heavy). Ship physical device."],
    ["App Service Migration Assistant", "Web app migration", "Assess IIS/.NET apps, migrate to App Service."],
    ["Azure Migrate: Server Migration", "VM migration", "VMware (agentless), Hyper-V, physical. Continuous replication."],
    ["Storage Migration Service", "File server migration", "Migrate Windows file servers to Azure Files or Windows VMs."],
], col_widths=[Inches(2.8), Inches(2.5), Inches(6.2)],
    font_size=13)

# Slide 90: Demo
s = new_content_slide("Segment 5 Demos", 5)
add_demo_slide(s, [
    "Hub-Spoke + Peering: Create hub VNet with Firewall, spoke VNets with peering and UDRs",
    "Private Endpoint + DNS: Create a Private Endpoint for a storage account with Private DNS Zone",
    "NSG + Firewall: Configure NSG micro-segmentation and Azure Firewall application rules",
    "Azure Migrate Assessment: Run a migration assessment for an on-premises environment",
])

# Slide 91: Exam Tips
s = new_exam_tip_slide("Segment 5 Exam Tips", 5)
add_bullet_list(s, [
    "ExpressRoute is NOT encrypted by default. If a question asks for encrypted private connectivity, add IPsec VPN over ER or use MACsec (Direct).",
    "Private Endpoint + Private DNS Zone = Zero Trust networking. Default recommendation for all data services.",
    "Azure Firewall for centralized egress control. UDRs force spoke traffic through the hub firewall.",
    "Front Door for global HTTP. App Gateway for regional HTTP. Load Balancer for regional non-HTTP. Traffic Manager for DNS-based global.",
    "SQL Managed Instance = migration sweet spot. Near 100% SQL Server compat with fully managed PaaS. First choice for SQL migration.",
    "Virtual WAN for scale (10+ branches). Manual hub-spoke for customization and smaller environments.",
], font_size=18, top=Inches(1.5), color=DARK_BLUE)

# Slide 92: Review
s = new_content_slide("Segment 5 Review Questions", 5)
add_review_questions(s, [
    ("Your company has 50 branch offices connecting to Azure. They need any-to-any connectivity between branches and Azure VNets with centralized security policies. What do you recommend?",
     "Azure Virtual WAN with Secured Virtual Hub (integrated Firewall Manager). Provides managed routing, automatic any-to-any connectivity, and centralized security at scale for 50+ branches."),
    ("An application in a spoke VNet needs to access Azure SQL Database without any traffic traversing the public internet. How do you design this?",
     "Create a Private Endpoint for Azure SQL in the spoke VNet (or hub). Configure Private DNS Zone (privatelink.database.windows.net). Disable public access on SQL. UDR routes through hub firewall if cross-spoke."),
    ("You are migrating 200 SQL Server databases from on-prem. Most use SQL Agent, cross-database queries, and CLR. You want minimal application changes. What is your migration target and tool?",
     "Target: Azure SQL Managed Instance (supports all three features). Tool: Azure Database Migration Service (DMS) in online mode for minimal downtime. Assess with Azure Migrate first."),
])


# ---------- SECTION 6: CLOSING (Slides 93-97) ----------

# Slide 93: Recap
s = new_content_slide("Course Recap", 6, "5 Segments, 3 Cross-Cutting Themes")
add_table(s, [
    ["Segment", "Core Theme", "Key Takeaway"],
    ["1. Identity & Governance", "Who & What", "Managed Identity + PHS + PIM + Policy"],
    ["2. Data Storage", "Store & Protect", "SQL MI + Cosmos DB + Private Endpoints"],
    ["3. BCDR", "Survive & Recover", "AZ + Failover Groups + ASR + Backup"],
    ["4. Compute & Apps", "Build & Run", "Container Apps + Service Bus + APIM"],
    ["5. Network & Migration", "Connect & Move", "Hub-Spoke + ExpressRoute + Azure Migrate"],
], col_widths=[Inches(2.5), Inches(2), Inches(7)])
_add_textbox(s, Inches(0.8), Inches(5.0), Inches(11), Inches(0.8),
             "Everywhere: Well-Architected Framework | Zero Trust | Managed Identity",
             font_size=20, bold=True, color=AZURE_BLUE)

# Slide 94: Exam Day Strategy
s = new_content_slide("Exam Day Strategy", 6)
add_bullet_list(s, [
    "Pace: ~2 minutes per question. 40-60 questions in 120 minutes. Do NOT spend 5+ minutes on one question.",
    "Case Studies: Read the questions FIRST, then scan the case study for relevant details. Do not read the entire case study top-to-bottom.",
    "Elimination: Remove obviously wrong answers first. Usually 1-2 are clearly wrong. Then compare remaining options.",
    "Flag & Return: Flag uncertain questions and return after completing all others. Fresh perspective helps.",
    "No Penalty: There is NO penalty for guessing. NEVER leave a question blank. Always select your best guess.",
    "Multiple Correct: Some questions have multiple correct answers but ask for the BEST answer. Look for 'most cost-effective', 'least effort', 'most secure'.",
    "Read Carefully: Watch for qualifiers: 'minimum', 'maximum', 'only', 'must', 'should'. These change the correct answer.",
], font_size=17)

# Slide 95: Study Resources
s = new_content_slide("Study Resources", 6)
add_bullet_list(s, [
    "Microsoft Learn: Free learning paths aligned to AZ-305 exam objectives",
    "AZ-305 Study Guide: https://aka.ms/AZ305-StudyGuide -- official skills measured document",
    "Practice Assessment: Free practice questions on Microsoft Learn -- take at least 3 times",
    "Exam Sandbox: https://aka.ms/examdemo -- experience the exam interface before test day",
    "Azure Architecture Center: https://learn.microsoft.com/azure/architecture/ -- reference architectures, patterns, best practices",
    "Well-Architected Framework: https://learn.microsoft.com/azure/well-architected/ -- 5 pillars deep dive",
    "GitHub Labs: https://github.com/MicrosoftLearning/AZ-305-DesigningMicrosoftAzureInfrastructureSolutions",
    "This Course Repo: Contains all demo scripts, Bicep templates, and practice questions",
], font_size=17)

# Slide 96: Next Steps
s = new_content_slide("Three Prioritized Next Steps", 6)
# Use colored boxes for each priority
priorities = [
    ("[IMMEDIATE]", "Take the Microsoft Learn practice assessment TODAY. Identify your weak areas immediately.",
     GREEN, Inches(1.5)),
    ("[SHORT-TERM]", "Build 3 architectures in an Azure sandbox THIS WEEK: hub-spoke networking, SQL failover group, Container Apps deployment.",
     AZURE_BLUE, Inches(3.2)),
    ("[LONG-TERM]", "Schedule the AZ-305 exam within 2 WEEKS. A deadline creates urgency. You have the knowledge -- go prove it.",
     PURPLE, Inches(4.9)),
]
for label, desc, color, top_pos in priorities:
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.8), top_pos,
                             Inches(11.5), Inches(1.3))
    _set_fill_solid(box, color)
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = label + "  "
    run1.font.size = Pt(20)
    run1.font.bold = True
    run1.font.color.rgb = GOLD
    run1.font.name = "Segoe UI Semibold"
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(18)
    run2.font.color.rgb = WHITE
    run2.font.name = "Segoe UI"


# Slide 97: Thank You
slide_counter += 1
s = prs.slides.add_slide(BLANK_LAYOUT)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                        SLIDE_WIDTH, SLIDE_HEIGHT)
_set_fill_solid(bg, DARK_BLUE)
bg.line.fill.background()

bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.3),
                         SLIDE_WIDTH, Inches(0.06))
_set_fill_solid(bar, GOLD)
bar.line.fill.background()

_add_textbox(s, Inches(1), Inches(1.0), Inches(11), Inches(1),
             "Thank You!", font_size=48, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold")

_add_textbox(s, Inches(1), Inches(2.0), Inches(11), Inches(1),
             "Questions? Let's discuss!", font_size=28, color=GOLD,
             alignment=PP_ALIGN.CENTER)

_add_textbox(s, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
             "Tim Warner", font_size=24, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)

_add_textbox(s, Inches(1), Inches(4.5), Inches(11), Inches(2),
             "timothywarner316@gmail.com  |  @TechTrainerTim\n"
             "O'Reilly Live Learning\n"
             "Microsoft MVP  |  MCT  |  Azure Solutions Architect Expert",
             font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)

add_footer(s, 6, slide_counter)


# ============================================================================
# SAVE
# ============================================================================
output_path = os.path.join(os.path.dirname(__file__),
                           "warner-az-305-2026-comprehensive.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {slide_counter}")
