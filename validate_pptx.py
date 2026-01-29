"""
AZ-305 PPTX Validation Script
Validates the comprehensive presentation for completeness, accuracy, and quality.
"""

import sys
from pptx import Presentation
from pptx.util import Inches

PPTX_PATH = r"C:\github\az305\warner-az-305-2026-comprehensive.pptx"

def extract_slide_info(prs):
    """Extract title and content summary from each slide."""
    slides_info = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(row_texts))
        # First substantial text is usually the title
        title = texts[0] if texts else "(no text)"
        slides_info.append({
            "num": i,
            "title": title,
            "all_text": texts,
            "full_text": "\n".join(texts),
        })
    return slides_info


def check_slide_count(slides_info):
    """Check that all 97 slides are present."""
    count = len(slides_info)
    status = "PASS" if count == 97 else "FAIL"
    return status, f"Slide count: {count} (expected 97)"


def check_segment_dividers(slides_info):
    """Check that 5 segment dividers exist with correct exam weights."""
    expected_dividers = [
        ("Identity, Governance", "25-30%"),
        ("Data Storage", "20-25%"),
        ("Business Continuity", "15-20%"),
        ("Compute", "~17%"),
        ("Networking", "~18%"),
    ]
    found = []
    issues = []
    for topic, weight in expected_dividers:
        matched = False
        for s in slides_info:
            full = s["full_text"].lower()
            if topic.lower() in full and weight in s["full_text"]:
                found.append((topic, s["num"]))
                matched = True
                break
        if not matched:
            issues.append(f"Missing divider for '{topic}' with weight '{weight}'")

    status = "PASS" if len(found) == 5 else "FAIL"
    detail = f"Found {len(found)}/5 segment dividers"
    if issues:
        detail += " | Issues: " + "; ".join(issues)
    return status, detail, found


def check_segment_structure(slides_info):
    """Check each segment has: learning objectives, content, demo, exam tips, review questions."""
    segments = {
        1: {"name": "Identity/Governance", "range": (9, 27)},
        2: {"name": "Data Storage", "range": (28, 45)},
        3: {"name": "BCDR", "range": (46, 60)},
        4: {"name": "Compute/App", "range": (61, 76)},
        5: {"name": "Networking/Migration", "range": (77, 92)},
    }
    results = []
    for seg_id, info in segments.items():
        start, end = info["range"]
        seg_slides = [s for s in slides_info if start <= s["num"] <= end]
        seg_text = " ".join([s["full_text"].lower() for s in seg_slides])

        checks = {
            "learning_objectives": any("learning objective" in s["full_text"].lower() for s in seg_slides),
            "demo_callout": any("demo" in s["full_text"].lower() for s in seg_slides),
            "exam_tips": any("exam tip" in s["full_text"].lower() for s in seg_slides),
            "review_questions": any("review question" in s["full_text"].lower() for s in seg_slides),
        }
        missing = [k for k, v in checks.items() if not v]
        status = "PASS" if not missing else "FAIL"
        detail = f"Segment {seg_id} ({info['name']}): "
        if missing:
            detail += f"MISSING: {', '.join(missing)}"
        else:
            detail += "All required sections present"
        results.append((status, detail))
    return results


def check_decision_matrices(slides_info):
    """Check for decision matrix tables on key topics."""
    key_matrices = [
        "authentication decision",
        "storage decision",
        "redundancy",
        "rto",
        "compute decision",
        "messaging",
        "load balancing",
        "migration",
        "firewall",
        "governance decision",
    ]
    found = []
    missing = []
    for matrix in key_matrices:
        matched = False
        for s in slides_info:
            if matrix.lower() in s["full_text"].lower():
                found.append(matrix)
                matched = True
                break
        if not matched:
            missing.append(matrix)

    status = "PASS" if len(found) >= 8 else "WARN" if len(found) >= 6 else "FAIL"
    detail = f"Decision matrices found: {len(found)}/{len(key_matrices)}"
    if missing:
        detail += f" | Not found: {', '.join(missing)}"
    return status, detail


def check_opening_closing(slides_info):
    """Check opening (8 slides) and closing (5 slides)."""
    opening_keywords = ["az-305", "about", "course objectives", "exam overview",
                        "agenda", "cross-cutting", "well-architected", "how to use"]
    closing_keywords = ["recap", "exam day", "study resources", "next steps", "thank you"]

    opening_slides = slides_info[:8]
    closing_slides = slides_info[-5:]

    opening_found = sum(1 for kw in opening_keywords
                        if any(kw.lower() in s["full_text"].lower() for s in opening_slides))
    closing_found = sum(1 for kw in closing_keywords
                        if any(kw.lower() in s["full_text"].lower() for s in closing_slides))

    o_status = "PASS" if opening_found >= 7 else "WARN"
    c_status = "PASS" if closing_found >= 4 else "WARN"
    return (o_status, f"Opening: {opening_found}/{len(opening_keywords)} keywords matched"),\
           (c_status, f"Closing: {closing_found}/{len(closing_keywords)} keywords matched")


def check_course_flow_alignment(slides_info):
    """Cross-reference with course flow topics."""
    # Key topics from course flow that should appear in the deck
    course_flow_topics = [
        # Segment 1
        "entra id", "conditional access", "rbac", "management group", "azure policy",
        "key vault", "azure monitor", "pim", "sentinel",
        # Segment 2
        "sql database", "sql managed instance", "cosmos db", "blob storage",
        "data lake", "data factory", "synapse", "partition key",
        # Segment 3
        "rto", "rpo", "availability zone", "availability set", "site recovery",
        "backup", "failover group", "sla",
        # Segment 4
        "container apps", "aks", "azure functions", "service bus",
        "event grid", "event hubs", "api management", "redis cache",
        # Segment 5
        "hub-spoke", "expressroute", "vpn gateway", "nsg", "azure firewall",
        "private endpoint", "front door", "azure migrate", "virtual wan",
    ]
    all_text = " ".join([s["full_text"].lower() for s in slides_info])
    found = [t for t in course_flow_topics if t.lower() in all_text]
    missing = [t for t in course_flow_topics if t.lower() not in all_text]

    status = "PASS" if len(missing) == 0 else "WARN" if len(missing) <= 3 else "FAIL"
    detail = f"Course flow topics covered: {len(found)}/{len(course_flow_topics)}"
    if missing:
        detail += f" | Missing: {', '.join(missing)}"
    return status, detail


def check_exam_weights(slides_info):
    """Verify exam weights are mentioned correctly."""
    weights = {
        "25-30%": "Identity/Governance",
        "20-25%": "Data Storage",
        "15-20%": "Business Continuity",
    }
    all_text = " ".join([s["full_text"] for s in slides_info])
    found = {w: name for w, name in weights.items() if w in all_text}
    status = "PASS" if len(found) == len(weights) else "WARN"
    return status, f"Exam weights verified: {len(found)}/{len(weights)}"


def check_tables_present(slides_info):
    """Count slides that contain table markers (pipe characters in text)."""
    table_count = 0
    for s in slides_info:
        if any(" | " in t for t in s["all_text"]):
            table_count += 1
    status = "PASS" if table_count >= 15 else "WARN"
    return status, f"Slides with tables: {table_count}"


def main():
    print("=" * 70)
    print("AZ-305 PPTX VALIDATION REPORT")
    print("=" * 70)

    # Load PPTX
    try:
        prs = Presentation(PPTX_PATH)
        print(f"\n[PASS] PPTX loaded successfully: {PPTX_PATH}")
        print(f"       Slide dimensions: {prs.slide_width} x {prs.slide_height}")
    except Exception as e:
        print(f"\n[FAIL] Could not load PPTX: {e}")
        sys.exit(1)

    slides_info = extract_slide_info(prs)

    print(f"\n{'='*70}")
    print("1. SLIDE COUNT")
    print(f"{'='*70}")
    status, detail = check_slide_count(slides_info)
    print(f"   [{status}] {detail}")

    print(f"\n{'='*70}")
    print("2. SEGMENT DIVIDERS")
    print(f"{'='*70}")
    status, detail, found = check_segment_dividers(slides_info)
    print(f"   [{status}] {detail}")
    for topic, num in found:
        print(f"     - Slide {num}: {topic}")

    print(f"\n{'='*70}")
    print("3. SEGMENT STRUCTURE (Learning Obj, Content, Demo, Exam Tips, Review)")
    print(f"{'='*70}")
    results = check_segment_structure(slides_info)
    for status, detail in results:
        print(f"   [{status}] {detail}")

    print(f"\n{'='*70}")
    print("4. DECISION MATRICES")
    print(f"{'='*70}")
    status, detail = check_decision_matrices(slides_info)
    print(f"   [{status}] {detail}")

    print(f"\n{'='*70}")
    print("5. OPENING & CLOSING SECTIONS")
    print(f"{'='*70}")
    (o_s, o_d), (c_s, c_d) = check_opening_closing(slides_info)
    print(f"   [{o_s}] {o_d}")
    print(f"   [{c_s}] {c_d}")

    print(f"\n{'='*70}")
    print("6. COURSE FLOW ALIGNMENT")
    print(f"{'='*70}")
    status, detail = check_course_flow_alignment(slides_info)
    print(f"   [{status}] {detail}")

    print(f"\n{'='*70}")
    print("7. EXAM WEIGHTS")
    print(f"{'='*70}")
    status, detail = check_exam_weights(slides_info)
    print(f"   [{status}] {detail}")

    print(f"\n{'='*70}")
    print("8. TABLE CONTENT DENSITY")
    print(f"{'='*70}")
    status, detail = check_tables_present(slides_info)
    print(f"   [{status}] {detail}")

    # Print slide manifest
    print(f"\n{'='*70}")
    print("SLIDE MANIFEST")
    print(f"{'='*70}")
    for s in slides_info:
        title_short = s["title"][:75]
        print(f"   Slide {s['num']:2d}: {title_short}")

    # Grade calculation
    print(f"\n{'='*70}")
    print("QUALITY ASSESSMENT")
    print(f"{'='*70}")

    # Collect all statuses
    all_checks = []
    all_checks.append(check_slide_count(slides_info)[0])
    all_checks.append(check_segment_dividers(slides_info)[0])
    for r in check_segment_structure(slides_info):
        all_checks.append(r[0])
    all_checks.append(check_decision_matrices(slides_info)[0])
    all_checks.append(check_opening_closing(slides_info)[0][0])
    all_checks.append(check_opening_closing(slides_info)[1][0])
    all_checks.append(check_course_flow_alignment(slides_info)[0])
    all_checks.append(check_exam_weights(slides_info)[0])
    all_checks.append(check_tables_present(slides_info)[0])

    passes = all_checks.count("PASS")
    warns = all_checks.count("WARN")
    fails = all_checks.count("FAIL")
    total = len(all_checks)

    print(f"   PASS: {passes}/{total}")
    print(f"   WARN: {warns}/{total}")
    print(f"   FAIL: {fails}/{total}")

    if fails == 0 and warns == 0:
        grade = "A"
    elif fails == 0 and warns <= 2:
        grade = "A-"
    elif fails == 0:
        grade = "B+"
    elif fails <= 1:
        grade = "B"
    elif fails <= 2:
        grade = "C"
    else:
        grade = "D"

    print(f"\n   OVERALL GRADE: {grade}")
    print(f"{'='*70}")


if __name__ == "__main__":
    main()
